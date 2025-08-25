from pptx import Presentation


def export_to_pptx(slides: list[dict], filename= "generated_slides.pptx"):
    p = Presentation()

    for slide_data in slides:
        slide_layout = p.slide_layouts[1]  # Title + Content
        slide = p.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data["title"]


        body = slide.placeholders[1].text_frame

        for bullet in slide_data["bullets"]:
            ps = body.add_paragraph()
            ps.text = bullet

    p.save(filename)