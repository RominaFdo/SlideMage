# from pptx import Presentation


# def export_to_pptx(slides: list[dict], filename= "generated_slides.pptx"):
#     p = Presentation()

#     for slide_data in slides:
#         slide_layout = p.slide_layouts[1]  # Title + Content
#         slide = p.slides.add_slide(slide_layout)
#         slide.shapes.title.text = slide_data["title"]


#         body = slide.placeholders[1].text_frame

#         for bullet in slide_data["bullets"]:
#             ps = body.add_paragraph()
#             ps.text = bullet

#     p.save(filename)
#     print(f"✅ Exported to {filename}")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import List, Dict, Any
import logging
import os

logger = logging.getLogger(__name__)

def export_to_pptx(slides: List[Dict[str, Any]], filename: str = "generated_slides.pptx") -> bool:
    try:
        prs = Presentation()
        
        if len(prs.slides) > 0:
            slides_to_remove = list(prs.slides)
            for slide in slides_to_remove:
                slide_id = slide.slide_id
                prs.part.drop_rel(prs.slides._sldIdLst[0].rId)
                del prs.slides._sldIdLst[0]
        
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        if slides:
            _create_title_slide(prs, slides[0]["title"])
        
        for i, slide_data in enumerate(slides):
            _create_content_slide(prs, slide_data, i + 1)
        
        prs.save(filename)
        logger.info(f"✅ Exported {len(slides)} slides to {filename}")
        return True
        
    except Exception as e:
        logger.error(f"❌ Export failed: {str(e)}")
        return False

def _create_title_slide(prs: Presentation, main_title: str):
    """Create an attractive title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = main_title.replace(" - Part 1", "").replace(" (Slide 1)", "")
        
        # Format title
        title_frame = title.text_frame
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Title font formatting
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.name = 'Calibri'
        title_run.font.size = Pt(44)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(31, 73, 125)  # Professional blue
    
    # Add subtitle
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "AI-Generated Presentation"
        
        # Format subtitle
        subtitle_frame = subtitle.text_frame
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        subtitle_run = subtitle_frame.paragraphs[0].runs[0]
        subtitle_run.font.name = 'Calibri'
        subtitle_run.font.size = Pt(24)
        subtitle_run.font.color.rgb = RGBColor(89, 89, 89)  # Gray

def _create_content_slide(prs: Presentation, slide_data: Dict[str, Any], slide_number: int):
    """Create a content slide with bullets."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = slide_data.get("title", f"Slide {slide_number}")
        
        # Format title
        title_frame = title.text_frame
        title_frame.margin_bottom = Inches(0.1)
        
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.name = 'Calibri'
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(31, 73, 125)
    
    # Add content
    if len(slide.placeholders) > 1:
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear default text
        
        bullets = slide_data.get("bullets", [])
        
        # Add bullets with formatting
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet_text
            p.level = 0  # First level bullet
            
            # Format bullet text
            if p.runs:
                run = p.runs[0]
                run.font.name = 'Calibri'
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(64, 64, 64)
        
        # Set text frame properties
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)

def create_enhanced_presentation(slides: List[Dict[str, Any]], filename: str = "presentation.pptx"):
    """
    Create a presentation with enhanced design elements.
    """
    try:
        prs = Presentation()
        
        # Customize slide master if needed
        _apply_theme(prs)
        
        # Create slides
        for i, slide_data in enumerate(slides):
            if i == 0:
                # Title slide
                _create_enhanced_title_slide(prs, slide_data)
            else:
                # Content slide
                _create_enhanced_content_slide(prs, slide_data, i)
        
        prs.save(filename)
        logger.info(f"✅ Created enhanced presentation: {filename}")
        return True
        
    except Exception as e:
        logger.error(f"❌ Enhanced presentation creation failed: {str(e)}")
        return False

def _apply_theme(prs: Presentation):
    """Apply a professional theme to the presentation."""
    # Note: Advanced theming would require more complex manipulation
    # of the slide master and layouts
    pass

def _create_enhanced_title_slide(prs: Presentation, slide_data: Dict[str, Any]):
    """Create an enhanced title slide with design elements."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Main title
    title = slide.shapes.title
    title.text = slide_data.get("title", "Presentation")
    
    # Enhanced title formatting
    title_frame = title.text_frame
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for run in title_frame.paragraphs[0].runs:
        run.font.name = 'Segoe UI'
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = RGBColor(31, 73, 125)
    
    # Add subtitle with first few bullets as overview
    if len(slide.placeholders) > 1 and slide_data.get("bullets"):
        subtitle = slide.placeholders[1]
        bullets = slide_data.get("bullets", [])
        subtitle_text = " • ".join(bullets[:3]) if len(bullets) > 1 else bullets[0] if bullets else ""
        subtitle.text = subtitle_text
        
        # Format subtitle
        subtitle_frame = subtitle.text_frame
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        for run in subtitle_frame.paragraphs[0].runs:
            run.font.name = 'Segoe UI'
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(89, 89, 89)

def _create_enhanced_content_slide(prs: Presentation, slide_data: Dict[str, Any], slide_number: int):
    """Create enhanced content slide with better formatting."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = slide_data.get("title", f"Content {slide_number}")
    
    # Enhanced title formatting
    title_frame = title.text_frame
    for run in title_frame.paragraphs[0].runs:
        run.font.name = 'Segoe UI'
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(31, 73, 125)
    
    # Content with enhanced bullets
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()
    
    bullets = slide_data.get("bullets", [])
    
    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet_text
        p.level = 0
        
        # Enhanced bullet formatting
        for run in p.runs:
            run.font.name = 'Segoe UI'
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(64, 64, 64)
        
        # Add spacing between bullets
        p.space_after = Pt(12)
    
    # Adjust text frame margins
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.2)